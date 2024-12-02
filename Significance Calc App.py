try:
    import customtkinter as ctk
except ImportError:
    print ("Please install pip3 install customtkinter to use this feature")
    exit()
from tkinter.colorchooser import askcolor  # For color picking
import os
import platform
import tkinter.filedialog as fd
from tkinter import messagebox
from scipy import stats
import math
import seaborn as sns
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
import matplotlib.pyplot as plt
import pandas as pd
from datetime import datetime #importing datetime to use for auto file creation
try:
    from pptx import Presentation  # Import for PowerPoint
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN  # Import for alignment
except ImportError:
    print("Please install python-pptx to use this feature")
    exit()
#from pptx import Presentation (imports for code that I am working on to make the app a bit more complex and output a power point slide. I am still working through the ways to make this function)
#from pptx.util import Inches

class SignificanceCalculatorApp:
    def __init__(self, master):
        self.master=master
        self.master.title("Statistical Significance Calculator")
        self.master.geometry("650x650")

        #setting a theme for the application
        ctk.set_appearance_mode("dark")
        ctk.set_default_color_theme("dark-blue")

        # Initialize UI components
        self.create_widgets()

    def create_widgets(self):
        # this is to prompt the input of a custom slide title.
        ctk.CTkLabel(self.master, text="Enter a Title for your slide:").grid(row=7, column=0)
        self.entry_slide_title = ctk.CTkEntry(self.master,width=150)
        self.entry_slide_title.grid(row=7, column=1)

        # Input fields
        ctk.CTkLabel(self.master, text="Sample Size A:").grid(row=0, column=0, padx=10, pady=5)
        self.entry_sample_size_a = ctk.CTkEntry(self.master)
        self.entry_sample_size_a.grid(row=0, column=1, padx=10, pady=5)

        ctk.CTkLabel(self.master, text="Percentage A:").grid(row=1, column=0, padx=10,pady=5)
        self.entry_percentage_a = ctk.CTkEntry(self.master)
        self.entry_percentage_a.grid(row=1, column=1, padx=10, pady=5)

        ctk.CTkLabel(self.master, text="Sample Size B:").grid(row=2, column=0, padx=10, pady=5)
        self.entry_sample_size_b = ctk.CTkEntry(self.master)
        self.entry_sample_size_b.grid(row=2, column=1, padx=10, pady=5)

        ctk.CTkLabel(self.master, text="Percentage B:").grid(row=3, column=0, padx=10, pady=5)
        self.entry_percentage_b = ctk.CTkEntry(self.master)
        self.entry_percentage_b.grid(row=3, column=1, padx=10, pady=5)

        # Color Picker for chart colors
        ctk.CTkButton(self.master, text="Pick Bar Color", command=self.pick_color).grid(row=4, column=0, pady=5)
       
        # Compute button
        compute_button = ctk.CTkButton(self.master, text="Compute", command=self.calculate_significance)
        compute_button.grid(row=4, column=0, columnspan=3, pady=5)

        # Reset button
        reset_button = ctk.CTkButton(self.master, text="Reset", command=self.reset_fields)
        reset_button.grid(row=4, column=2, padx=5, pady=5)

        # Output label
        self.output_text = ctk.StringVar()
        output_label = ctk.CTkLabel(self.master, textvariable=self.output_text)
        output_label.grid(row=5, column=1, columnspan=2)

        #output button to export graph
        export_button=ctk.CTkButton(self.master, text="Export to Power Point", command=self.export_to_powerpoint)
        export_button.grid(row=8, column=0, columnspan=3,padx=5, pady=5)

        # Frame for the plot
        self.plot_frame = ctk.CTkFrame(self.master)
        self.plot_frame.grid(row=6, column=0, columnspan=3, padx=20, pady=10)

    #This input def allows for validation of the percentages.
    def validate_input(self):
        try:
            # Get and validate inputs
            n1 = int(self.entry_sample_size_a.get())
            n2 = int(self.entry_sample_size_b.get())
            p1 = float(self.entry_percentage_a.get()) / 100
            p2 = float(self.entry_percentage_b.get()) / 100
            
            # Check if percentages are within the valid range
            if not (0 <= p1 <= 1 and 0 <= p2 <= 1):
                raise ValueError("Percentages should be between 0 and 100")
            
            return n1, n2, p1, p2
        
        except ValueError as e:
            # Display the error message
            self.output_text.set(f"Invalid input: {e}")
            return None

    def calculate_significance(self):
        # Validate inputs first
        inputs = self.validate_input()
        if inputs is None:
            return  # Exit if inputs are invalid

        # Unpack validated inputs
        n1, n2, p1, p2 = inputs

        # Calculate pooled proportion
        p_pool = (p1 * n1 + p2 * n2) / (n1 + n2)
        
        # Calculate the z-score
        z = (p1 - p2) / math.sqrt(p_pool * (1 - p_pool) * (1/n1 + 1/n2))
        
        # Calculate p-value from z-score (two-tailed test)
        p_value = stats.norm.sf(abs(z)) * 2
        
        # Variables to hold the significance level and status
        confidence_reached = None
        
        # Loop through different significance levels (80%-99%). I ordered descending here as you really want to understand the highest level it is significant at.
        for confidence in [0.99, 0.95, 0.90, 0.85, 0.80]:
            significance_level = 1 - confidence
            if p_value < significance_level:
                confidence_reached = confidence
                self.output_text.set(f"Significant at {confidence*100}% level (p-value: {p_value:.4f})")
                break
        else:
            self.output_text.set(f"Not significant (p-value: {p_value:.4f})")

    # Call the function to update the graph
        self.update_graph(n1,n2,p1, p2, confidence_reached)

    def pick_color(self):
        """Opens a color picker dialog to select two colors for the bars."""
        color1 = askcolor()[1]  # Get the hex color code for the first bar
        color2 = askcolor()[1]  # Get the hex color code for the second bar
        if color1 and color2:  # Ensure both colors are selected
            self.bar_colors = [color1, color2]  # Store the selected colors
        else:
            self.bar_colors = ['#1f77b4', '#ff7f0e']  # Default colors (blue and orange)

    def update_graph(self, n1, n2, p1, p2, confidence_reached):

        # Clear the previous plot
        for widget in self.plot_frame.winfo_children():
            widget.destroy()

        # Create a new figure
        fig, ax = plt.subplots(figsize=(6, 3))

        #setting themes
        sns.set_theme(style="whitegrid")
        sns.set_palette("pastel")
        
        # Data for the bar chart
        plt.xlabel(' ', fontsize=14)
        labels = ['Percentage A', 'Percentage B']
        values = [p1 * 100, p2 * 100]

        # Create a DataFrame for easier plotting with seaborn
        data = pd.DataFrame({'Labels': labels, 'Values': values})

        # Use the selected bar colors (or default ones)
        bar_colors = getattr(self, 'bar_colors', ['#1f77b4', '#ff7f0e'])  # Default colors if not selected

        # Use seaborn to create the bar plot
        sns.barplot(x='Labels', y='Values', data=data, ax=ax, palette=bar_colors)
        for index, value in enumerate(values):
         ax.text(index, value + 1, f"{value:.1f}%", ha='center', va='bottom', fontsize=10)

        # Set y-axis to always show 0 to 100
        ax.set_ylim(0, 100)

        # Add a line indicating the significance level, if any
        if confidence_reached:
            ax.axhline(confidence_reached * 100, color='red', linestyle='--', label=f'Significant at {confidence_reached*100}%')
            ax.legend()

        # Set the plot labels and title
        ax.set_ylabel('Percentages')
        ax.set_title('Percentage Comparison')

        #Calculate actual difference between percentages
        actual_difference= abs(p1-p2)*100

       # Define significance levels and required differences
        significance_levels = [0.80, 0.85, 0.90, 0.95, 0.99]
        required_differences = [
            round(stats.norm.ppf(1 - (1 - conf) / 2) * math.sqrt(p1 * (1 - p1) / n1 + p2 * (1 - p2) / n2), 2)
            for conf in significance_levels
        ]

        # Create table data
        table_data = [["Level", "Required Difference"]] + [
            [f"{int(conf * 100)}%", f"{diff:.2f}"] for conf, diff in zip(significance_levels, required_differences)
        ]

        # Create the table
        table = ax.table(cellText=table_data, colWidths=[0.3, 0.5], loc='right', cellLoc='center')
        table.auto_set_font_size(False)
        table.set_fontsize(8)
        table.scale(1, 1.5)  # Adjust table scale

        # Highlight the row based on confidence_reached
        if confidence_reached:
            # Find the index of the confidence_reached in significance_levels
            highlight_row_index = significance_levels.index(confidence_reached)

            # Highlight the corresponding row in the table
            for col in range(len(table_data[0])):
                cell = table[highlight_row_index + 1, col]  # +1 accounts for the header row
                cell.set_facecolor('#FFDDC1')  # Highlight color
                cell.set_fontsize(9)

        # Adjust the plot layout to accommodate the table
        plt.subplots_adjust(right=0.6)

        # Embed the plot into Tkinter
        canvas = FigureCanvasTkAgg(fig, master=self.plot_frame)
        canvas.draw()
        canvas.get_tk_widget().pack()

        #saving the plot as an image for export.
        fig.savefig('plot.png', bbox_inches='tight')


    def reset_fields(self): # This enables a reset button to clear fields and get ready for the next computation
        self.entry_sample_size_a.delete(0, ctk.END)
        self.entry_percentage_a.delete(0, ctk.END)
        self.entry_sample_size_b.delete(0, ctk.END)  # Fixed the reference here
        self.entry_percentage_b.delete(0, ctk.END)  # Fixed the reference here
        self.output_text.set('')  # Clear output text after reset

    # Clear the graph
        for widget in self.plot_frame.winfo_children():
            widget.destroy()

    def get_templates_directory():
            if platform.system() == "Darwin":  # macOS
                return os.path.expanduser("~/Library/Application Support/Microsoft/Office/User Templates/My Templates")
            elif platform.system() == "Windows":  # Windows
                return os.path.expandvars(r"%AppData%\Microsoft\Templates")
            else:
                return None  # For other OS, leave it blank or handle accordingly
    
    def export_to_powerpoint(self):
     
    # Prompt user to select a PowerPoint template
        template_path = fd.askopenfilename(
        title="Select PowerPoint To Append Slide To. If no template hit cancel for blank template",
        filetypes=[("PowerPoint files", "*.pptx")],
    )

        if not template_path:
        # If no template is selected, notify the user
            response = messagebox.askyesno(
            "No Template Selected",
            "No template was selected. Do you want to proceed with a blank presentation?",
        )
            if not response:
                return  # Exit if the user doesn't want to proceed

    # Load the template or create a new presentation
        try:
            if template_path:
                prs = Presentation(template_path)  # Load selected template
            else:
                prs = Presentation()  # Create a blank presentation
        except Exception as e:
            self.output_text.set(f"Error loading template: {e}")
            return

        slide_title = self.entry_slide_title.get()
        if not slide_title:
            self.output_text.set("Please provide a title for your slide.")
            return

        # Add a slide (assuming title layout exists in the template)
        try:
            slide_layout = prs.slide_layouts[6]  # Default to blank layout
        except IndexError:
            slide_layout = prs.slide_layouts[0]  # Fallback to title layout
        slide = prs.slides.add_slide(slide_layout)

        # Add a title as a text box
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(8)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = slide_title
        for paragraph in title_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Inches(0.375)

        # Add the plot to the slide
        img_path = getattr(self, 'latest_image_path', 'plot.png')
        left = Inches(0.1)  # Position from the left
        top = Inches(1.8)   # Position from the top
        width = Inches(9)   # Change the width as desired
        height = Inches(4.5)  # Change the height as desired

        # Add the picture to the slide with adjusted size
        slide.shapes.add_picture(img_path, left, top, width=width, height=height)
        
        # Add the significance output to the slide
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(8), Inches(1))
        text_frame = textbox.text_frame
        text_frame.text = self.output_text.get()

        # Save the PowerPoint file
        try:
            save_path = fd.asksaveasfilename(
                defaultextension=".pptx",
                filetypes=[("PowerPoint files", "*.pptx")],
                title="Save PowerPoint File",
            )
            if not save_path:
                self.output_text.set("Export canceled.")
                return

            prs.save(save_path)
            self.output_text.set(f"Exported to {save_path}")
        except Exception as e:
            self.output_text.set(f"Export failed: {e}")

# Tkinter setup
if __name__ == "__main__":
    root = ctk.CTk()
    app = SignificanceCalculatorApp(root)
    root.mainloop()
